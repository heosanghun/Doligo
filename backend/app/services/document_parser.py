"""Document parser for PDF/PPT/HWP."""
from pathlib import Path
import re


def parse_document(file_path: Path) -> dict | None:
    """Parse document and extract structure (sections, titles)."""
    ext = file_path.suffix.lower()
    if ext == ".pdf":
        return _parse_pdf(file_path)
    if ext == ".pptx":
        return _parse_pptx(file_path)
    if ext in (".hwp", ".hwpx"):
        return _parse_hwp(file_path)
    return None


def _parse_pdf(file_path: Path) -> dict | None:
    """Parse PDF structure."""
    from pypdf import PdfReader

    try:
        reader = PdfReader(str(file_path))
        full_text = ""
        for page in reader.pages:
            full_text += page.extract_text() or ""

        sections = _extract_sections_from_text(full_text)
        if not sections:
            # Fallback: one section per page
            sections = [{"title": f"페이지 {i+1}", "has_image_placeholder": True} for i in range(len(reader.pages))]

        return {
            "format": "pdf",
            "sections": sections,
            "raw_text_preview": full_text[:1000],
            "page_count": len(reader.pages),
        }
    except Exception:
        return None


def _extract_sections_from_text(text: str) -> list[dict]:
    """Extract sections from text using heuristics."""
    sections = []
    # Look for numbered headings: 1. Title, 2. Title, etc.
    pattern = r"(\d+)[\.\s]+([^\n]+)"
    matches = list(re.finditer(pattern, text))
    for m in matches:
        num, title = m.group(1), m.group(2).strip()
        if len(title) > 2 and len(title) < 100:
            sections.append({"title": title, "has_image_placeholder": True})
    if sections:
        return sections
    # Look for lines that look like headings (short, possibly bold indicators)
    lines = [l.strip() for l in text.split("\n") if l.strip()]
    for line in lines:
        if 3 <= len(line) <= 80 and not line.endswith("."):
            sections.append({"title": line, "has_image_placeholder": True})
            if len(sections) >= 10:
                break
    return sections[:10] if sections else []


def _parse_pptx(file_path: Path) -> dict | None:
    """Parse PowerPoint structure."""
    try:
        from pptx import Presentation

        prs = Presentation(str(file_path))
        sections = []
        for i, slide in enumerate(prs.slides):
            title = ""
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    title = shape.text[:100].strip() or f"슬라이드 {i+1}"
                    break
            if not title:
                title = f"슬라이드 {i+1}"
            sections.append({"title": title, "has_image_placeholder": True})
        return {
            "format": "ppt",
            "sections": sections,
            "raw_text_preview": "",
            "slide_count": len(prs.slides),
        }
    except Exception:
        return None


def _parse_hwp(file_path: Path) -> dict | None:
    """Parse HWP/HWPX structure (basic fallback)."""
    try:
        from hwp_hwpx_parser import Reader

        with Reader(str(file_path)) as r:
            text = r.text or ""
        sections = _extract_sections_from_text(text) if text else []
        if not sections:
            sections = [{"title": "본문", "has_image_placeholder": True}]
        return {
            "format": "hwp",
            "sections": sections,
            "raw_text_preview": text[:1000],
        }
    except ImportError:
        return {"format": "hwp", "sections": [{"title": "본문", "has_image_placeholder": True}], "raw_text_preview": ""}
    except Exception:
        return None
