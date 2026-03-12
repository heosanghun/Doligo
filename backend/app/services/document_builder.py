"""Document builder - fill template with content and images."""
import io
import logging
from pathlib import Path

from app.config import settings

logger = logging.getLogger(__name__)

# 한글 폰트 경로 후보 (우선순위 순)
_KOREAN_FONT_PATHS = [
    Path(__file__).resolve().parent.parent.parent / "fonts" / "NotoSansKR-Regular.ttf",
    Path(__file__).resolve().parent.parent.parent / "fonts" / "NanumGothic.ttf",
    Path("C:/Windows/Fonts/malgun.ttf"),  # Windows 맑은 고딕
    Path("C:/Windows/Fonts/gulim.ttc"),   # Windows 굴림
    Path("/usr/share/fonts/truetype/noto/NotoSansCJK-Regular.ttc"),
    Path("/usr/share/fonts/truetype/nanum/NanumGothic.ttf"),
]


def _get_korean_font_name() -> str:
    """한글 지원 TTF 폰트를 등록하고 폰트 이름 반환. 없으면 Helvetica."""
    from reportlab.pdfbase import pdfmetrics
    from reportlab.pdfbase.ttfonts import TTFont

    if getattr(_get_korean_font_name, "_registered", None):
        return "KoreanFont"

    for font_path in _KOREAN_FONT_PATHS:
        if font_path.exists():
            try:
                pdfmetrics.registerFont(TTFont("KoreanFont", str(font_path)))
                _get_korean_font_name._registered = True
                logger.info("한글 폰트 등록: %s", font_path)
                return "KoreanFont"
            except Exception as e:
                logger.warning("폰트 등록 실패 %s: %s", font_path, e)
    logger.warning("한글 폰트를 찾을 수 없음. Helvetica 사용 (한글 깨짐 가능)")
    return "Helvetica"


def build_document(
    template_path: Path,
    content: dict,
    structure: dict,
    output_format: str | None = None,
) -> Path:
    """Build final document. output_format: pdf, ppt, hwp, md. None이면 입력과 동일."""
    ext = template_path.suffix.lower()
    fmt = (output_format or "").lower() or _ext_to_format(ext)
    if fmt not in ("pdf", "ppt", "hwp", "md"):
        fmt = "pdf"

    # 출력 형식에 맞춰 빌드
    if fmt == "md":
        return _build_md_from_scratch(template_path, content, structure)
    if fmt == "pdf":
        if ext == ".pdf":
            return _build_pdf(template_path, content, structure)
        return _build_pdf_from_scratch(template_path, content, structure)
    if fmt == "ppt":
        if ext == ".pptx":
            return _build_pptx(template_path, content, structure)
        return _build_pptx_from_scratch(template_path, content, structure)
    if fmt == "hwp":
        return _build_hwp_from_scratch(template_path, content, structure)
    return _build_pdf(template_path, content, structure)


def _build_md_from_scratch(template_path: Path, content: dict, structure: dict) -> Path:
    """Markdown 형식으로 저장 - 복붙·편집에 최적."""
    sections = content.get("sections", [])
    lines = []
    for sec in sections:
        title = sec.get("title", "")
        text = sec.get("content", "").strip()
        if title:
            lines.append(f"# {title}\n")
        if text:
            lines.append(text)
            lines.append("\n\n")
    output_dir = Path(settings.output_dir)
    output_dir.mkdir(parents=True, exist_ok=True)
    out_path = output_dir / f"output_{template_path.stem}.md"
    out_path.write_text("".join(lines).strip(), encoding="utf-8")
    return out_path


def _ext_to_format(ext: str) -> str:
    if ext == ".pdf":
        return "pdf"
    if ext == ".pptx":
        return "ppt"
    if ext in (".hwp", ".hwpx"):
        return "hwp"
    return "pdf"


def _build_pdf(template_path: Path, content: dict, structure: dict) -> Path:
    """Build PDF with overlay (한글 폰트 + Paragraph 지원)."""
    from pypdf import PdfReader, PdfWriter
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.styles import ParagraphStyle
    from reportlab.lib.units import mm
    from reportlab.pdfgen import canvas
    from reportlab.platypus import Paragraph

    reader = PdfReader(str(template_path))
    writer = PdfWriter()
    sections = content.get("sections", [])
    images = content.get("images", {})
    font_name = _get_korean_font_name()
    width, height = A4

    style = ParagraphStyle(
        name="KoreanBody",
        fontName=font_name,
        fontSize=10,
        leading=14,
        wordWrap="CJK",
    )
    style_title = ParagraphStyle(
        name="KoreanTitle",
        fontName=font_name,
        fontSize=12,
        leading=16,
        spaceAfter=6,
    )

    for page_idx, page in enumerate(reader.pages):
        packet = io.BytesIO()
        c = canvas.Canvas(packet, pagesize=A4)

        sec_idx = min(page_idx, len(sections) - 1)
        if sec_idx >= 0 and sections:
            sec = sections[sec_idx]
            title = sec.get("title", "")
            text = sec.get("content", "")[:3000]
            y = height - 50

            if title:
                try:
                    p_title = Paragraph(title, style_title)
                    p_title.wrapOn(c, width - 100, 50)
                    p_title.drawOn(c, 50, y - 20)
                    y -= 30
                except Exception:
                    c.setFont(font_name, 12)
                    c.drawString(50, y - 20, title[:80])
                    y -= 25

            for line in text.split("\n")[:50]:
                line = line.strip()
                if not line:
                    y -= 8
                    continue
                try:
                    para = Paragraph(line.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;"), style)
                    para.wrapOn(c, width - 100, 30)
                    h = para.height
                    if y - h < 50:
                        break
                    para.drawOn(c, 50, y - h)
                    y -= h + 4
                except Exception:
                    c.setFont(font_name, 10)
                    c.drawString(50, y, line[:120])
                    y -= 14
                if y < 50:
                    break

        if sec_idx in images and images.get(sec_idx):
            try:
                from reportlab.lib.utils import ImageReader
                img = ImageReader(io.BytesIO(images[sec_idx]))
                c.drawImage(img, 50, height - 300, width=200, height=150)
            except Exception:
                pass

        c.save()
        packet.seek(0)
        from pypdf import PdfReader as PdfR
        overlay = PdfR(packet)
        page.merge_page(overlay.pages[0])
        writer.add_page(page)

    output_dir = Path(settings.output_dir)
    output_dir.mkdir(parents=True, exist_ok=True)
    out_path = output_dir / f"output_{template_path.stem}.pdf"
    with open(out_path, "wb") as f:
        writer.write(f)
    return out_path


def _build_pdf_from_scratch(template_path: Path, content: dict, structure: dict) -> Path:
    """Create PDF from content only (no template overlay)."""
    from pypdf import PdfWriter
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.styles import ParagraphStyle
    from reportlab.pdfgen import canvas
    from reportlab.platypus import Paragraph

    writer = PdfWriter()
    sections = content.get("sections", [])
    images = content.get("images", {})
    font_name = _get_korean_font_name()
    width, height = A4

    style = ParagraphStyle(
        name="KoreanBody",
        fontName=font_name,
        fontSize=10,
        leading=14,
        wordWrap="CJK",
    )
    style_title = ParagraphStyle(
        name="KoreanTitle",
        fontName=font_name,
        fontSize=12,
        leading=16,
        spaceAfter=6,
    )

    for sec_idx, sec in enumerate(sections):
        packet = io.BytesIO()
        c = canvas.Canvas(packet, pagesize=A4)
        title = sec.get("title", "")
        text = sec.get("content", "")[:3000]
        y = height - 50

        if title:
            try:
                p_title = Paragraph(title, style_title)
                p_title.wrapOn(c, width - 100, 50)
                p_title.drawOn(c, 50, y - 20)
                y -= 30
            except Exception:
                c.setFont(font_name, 12)
                c.drawString(50, y - 20, title[:80])
                y -= 25

        for line in text.split("\n")[:50]:
            line = line.strip()
            if not line:
                y -= 8
                continue
            try:
                para = Paragraph(line.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;"), style)
                para.wrapOn(c, width - 100, 30)
                h = para.height
                if y - h < 50:
                    break
                para.drawOn(c, 50, y - h)
                y -= h + 4
            except Exception:
                c.setFont(font_name, 10)
                c.drawString(50, y, line[:120])
                y -= 14
            if y < 50:
                break

        if sec_idx in images and images.get(sec_idx):
            try:
                from reportlab.lib.utils import ImageReader
                img = ImageReader(io.BytesIO(images[sec_idx]))
                c.drawImage(img, 50, height - 300, width=200, height=150)
            except Exception:
                pass

        c.save()
        packet.seek(0)
        from pypdf import PdfReader as PdfR
        page = PdfR(packet).pages[0]
        writer.add_page(page)

    output_dir = Path(settings.output_dir)
    output_dir.mkdir(parents=True, exist_ok=True)
    out_path = output_dir / f"output_{template_path.stem}.pdf"
    with open(out_path, "wb") as f:
        writer.write(f)
    return out_path


def _build_pptx_from_scratch(template_path: Path, content: dict, structure: dict) -> Path:
    """Create PPT from content only (new presentation)."""
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt

        prs = Presentation()
        sections = content.get("sections", [])
        images = content.get("images", {})

        for sec_idx, sec in enumerate(sections):
            layout = prs.slide_layouts[1]  # Title and Content
            slide = prs.slides.add_slide(layout)
            title = sec.get("title", "") or f"섹션 {sec_idx + 1}"
            text = sec.get("content", "")[:1000]
            slide.shapes.title.text = title
            body = slide.placeholders[1]
            tf = body.text_frame
            tf.clear()
            p = tf.paragraphs[0] if tf.paragraphs else tf.add_paragraph()
            p.text = text
            p.font.size = Pt(12)
            if sec_idx in images and images.get(sec_idx):
                try:
                    slide.shapes.add_picture(io.BytesIO(images[sec_idx]), Inches(0.5), Inches(4), width=Inches(3), height=Inches(2))
                except Exception:
                    pass

        output_dir = Path(settings.output_dir)
        output_dir.mkdir(parents=True, exist_ok=True)
        out_path = output_dir / f"output_{template_path.stem}.pptx"
        prs.save(str(out_path))
        return out_path
    except Exception as e:
        logger.warning("PPT 생성 실패: %s. PDF로 대체", e)
        return _build_pdf_from_scratch(template_path, content, structure)


def _build_hwp_from_scratch(template_path: Path, content: dict, structure: dict) -> Path:
    """Create HWPX from content. python-hwpx 필요."""
    try:
        from hwpx import HwpxDocument

        doc = HwpxDocument.new()
        sections = content.get("sections", [])
        for sec in sections:
            doc.add_paragraph(sec.get("title", ""))
            doc.add_paragraph(sec.get("content", "")[:2000])
        output_dir = Path(settings.output_dir)
        output_dir.mkdir(parents=True, exist_ok=True)
        out_path = output_dir / f"output_{template_path.stem}.hwpx"
        doc.save_to_path(str(out_path))
        return out_path
    except ImportError:
        logger.warning("python-hwpx 미설치. pip install python-hwpx 후 HWP 사용 가능. PDF로 대체 출력")
        return _build_pdf_from_scratch(template_path, content, structure)
    except Exception as e:
        logger.warning("HWP 생성 실패: %s. PDF로 대체", e)
        return _build_pdf_from_scratch(template_path, content, structure)


def _build_pptx(template_path: Path, content: dict, structure: dict) -> Path:
    """Build PPTX with content and images."""
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt

        prs = Presentation(str(template_path))
        sections = content.get("sections", [])
        images = content.get("images", {})

        for slide_idx, slide in enumerate(prs.slides):
            sec_idx = min(slide_idx, len(sections) - 1)
            if sec_idx < 0:
                continue
            sec = sections[sec_idx]
            for shape in slide.shapes:
                if hasattr(shape, "text_frame") and shape.text_frame:
                    shape.text_frame.clear()
                    p = shape.text_frame.paragraphs[0] if shape.text_frame.paragraphs else shape.text_frame.add_paragraph()
                    p.text = sec.get("content", "")[:500]
                    p.font.size = Pt(12)
                if shape.shape_type == 13:  # Picture placeholder
                    if sec_idx in images and images[sec_idx]:
                        try:
                            shape.insert_picture(io.BytesIO(images[sec_idx]))
                        except Exception:
                            pass

        output_dir = Path(settings.output_dir)
        output_dir.mkdir(parents=True, exist_ok=True)
        out_path = output_dir / f"output_{template_path.stem}.pptx"
        prs.save(str(out_path))
        return out_path
    except Exception:
        output_dir = Path(settings.output_dir)
        output_dir.mkdir(parents=True, exist_ok=True)
        out_path = output_dir / f"output_{template_path.stem}.pptx"
        out_path.write_bytes(template_path.read_bytes())
        return out_path
